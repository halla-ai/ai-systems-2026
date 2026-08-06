from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    def add_slide(title, points):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        
        tf = slide.placeholders[1].text_frame
        tf.word_wrap = True
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            if ":" in point:
                p.font.bold = True

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "TraceLedger: AI 협업 컨텍스트 자산화 엔진"
    subtitle.text = "휘발되는 채팅 로그로부터 엔지니어링 지식(Why) 추출 자동화\n\n발표자: [본인 성함]"

    # Slide 2: 배경 및 문제 정의
    add_slide("1. 배경: 지식의 휘발 (Problem)", [
        "현상: AI 협업 과정의 시행착오와 의사결정 근거가 종료 후 소멸",
        "문제: 코드 결과물만 남고 결정의 맥락은 블랙박스화",
        "해결: 대화 로그 분석 기반의 전문 기술 문서(Runbook, ADR) 자동화"
    ])

    # Slide 3: 시스템 아키텍처
    add_slide("2. 시스템 아키텍처 (L1-L7 Hierarchy)", [
        "L1-L2 (Infrastructure): 로그 전처리(Noise Filter) 및 토큰 실시간 추적",
        "L3 (Collaboration): Planner-Worker-Reviewer 기반 3인 체계",
        "L6 (Governance): 시도 횟수 제한 및 예산 기반 실행 통제",
        "L7 (Schema): Task Packet 규격을 통한 데이터 정합성 보장"
    ])

    # Slide 4: 핵심 워크플로우
    add_slide("3. 핵심 워크플로우: Ralph Loop", [
        "Planner (분석): Root Cause, 트레이드오프 분석",
        "Worker (작성): 분석 데이터 기반 Markdown 문서 생성",
        "Reviewer (검증): 원본 로그 대조 및 할루시네이션 검증",
        "Self-Correction: 반려 시 피드백 반영을 통한 즉각적 수정 루프"
    ])

    # Slide 5: 사례 분석
    add_slide("4. 사례 분석 (Case Study: React 인증)", [
        "입력: useEffect 무한 루프 -> useApi 훅 설계 -> Refresh Token 구현 로그",
        "Root Cause: 의존성 배열 상태 업데이트에 따른 재귀적 렌더링 분석",
        "Architecture: API 통신 레이어 표준화 및 커스텀 훅 설계 근거",
        "ADR: 401 에러 자동 갱신 처리에 대한 아키텍처 결정 기록"
    ])

    # Slide 6: 최적화 전략 (FinOps)
    add_slide("5. 최적화 전략: FinOps (비용 절감)", [
        "성과: 건당 약 40,000 토큰 -> 약 20,000 토큰 (50% 절감)",
        "Model Tiering: 분석(Planner)은 고사양, 작성/검증은 가성비 모델 배치",
        "Noise Filtering: 불필요한 인사말 및 메타데이터 필터링",
        "Incremental Analysis: 이전 분석 요약본 활용을 통한 중복 입력 최소화"
    ])

    # Slide 7: 운영 안정성 (Governance)
    add_slide("6. 운영 안정성: Smart Governance", [
        "Smart Compromise: 검증 점수 0.85 이상 시 즉시 승인 처리",
        "Final Fallback: 최대 시도 도달 시 '에이전트 제언' 포함 최종 저장",
        "결과: 무한 루프 방지 및 안정적인 결과물 반환 보장"
    ])

    # Slide 8: 결론
    add_slide("7. 결론 (Conclusion)", [
        "요약: 채팅 데이터를 팀의 영구적인 기술 자산으로 전환하는 엔진 구현",
        "효과: 코드 이해 시간 단축 및 엔지니어링 투명성 확보",
        "향후 계획: GitHub Action 연동을 통한 PR 기반 자동 분석 확장"
    ])

    prs.save('TraceLedger_Presentation.pptx')
    print("PPT generated successfully: TraceLedger_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
